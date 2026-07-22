"""
Build Everde AI Operations portal overview deck for IT/security meetings.
Output: docs/presentations/Everde_Portal_Overview_Jason_Aaron.pptx

  python scripts/presentations/build_portal_meeting_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "docs" / "presentations" / "Everde_Portal_Overview_Jason_Aaron.pptx"

FOREST = RGBColor(0x1B, 0x43, 0x32)
FOREST_DEEP = RGBColor(0x0F, 0x24, 0x1A)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CANVAS = RGBColor(0xF4, 0xF1, 0xEA)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
LAN_FILL = RGBColor(0xE8, 0xF0, 0xEB)
CLOUD_FILL = RGBColor(0xE3, 0xED, 0xF7)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_block(slide, title: str, subtitle: str = "") -> None:
    set_slide_bg(slide, FOREST_DEEP)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GOLD
        p2.space_before = Pt(12)


def add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, CANVAS)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(8.7), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16 if not line.startswith("  ") else 14)
        p.font.color.rgb = TEXT if not line.startswith("  ") else MUTED
        p.space_after = Pt(8)
        if line.startswith("•") or line.startswith("  •"):
            p.level = 1 if line.startswith("  ") else 0


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = WHITE

    cols, rws = len(headers), len(rows) + 1
    table = slide.shapes.add_table(rws, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * rws)).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = FOREST
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = "How data flows (two systems)"
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = WHITE

    # LAN box
    lan = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(4.2), Inches(4.5)
    )
    lan.fill.solid()
    lan.fill.fore_color.rgb = LAN_FILL
    lan.line.color.rgb = FOREST

    lan_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(3.9), Inches(0.4))
    p = lan_title.text_frame.paragraphs[0]
    p.text = "On-prem / VPN (Everde network)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = FOREST

    items = [
        "DataDrops share\n(Excel, xlsb, feeds)",
        "Agent PC\n(VRD-8FQJYW3)",
        "Python pipelines\n(freight, retail, sales plan, weather)",
    ]
    y = 1.85
    for label in items:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y), Inches(3.6), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = FOREST
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.12), Inches(3.3), Inches(0.7))
        tb.text_frame.paragraphs[0].text = label
        tb.text_frame.paragraphs[0].font.size = Pt(12)
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += 1.05

    # Cloud box
    cloud = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(1.25), Inches(4.2), Inches(4.5)
    )
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = CLOUD_FILL
    cloud.line.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    cloud_title = slide.shapes.add_textbox(Inches(5.5), Inches(1.35), Inches(3.9), Inches(0.4))
    cloud_title.text_frame.paragraphs[0].text = "Internet-facing (hosted)"
    cloud_title.text_frame.paragraphs[0].font.bold = True
    cloud_title.text_frame.paragraphs[0].font.size = Pt(14)
    cloud_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    cloud_items = [
        "Azure Blob Storage\n(JSON snapshots only)",
        "Vercel\nNext.js portal",
        "Users @everde.com\n(HTTPS + sign-in)",
    ]
    y = 1.85
    for label in cloud_items:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.65), Inches(y), Inches(3.6), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
        tb = slide.shapes.add_textbox(Inches(5.8), Inches(y + 0.12), Inches(3.3), Inches(0.7))
        tb.text_frame.paragraphs[0].text = label
        tb.text_frame.paragraphs[0].font.size = Pt(12)
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += 1.05

    # Arrows / labels between
    arrow1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(4.35), Inches(2.35), Inches(0.85), Inches(0.35)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = GOLD
    arrow1.line.fill.background()

    lbl1 = slide.shapes.add_textbox(Inches(4.0), Inches(2.75), Inches(1.4), Inches(0.35))
    lbl1.text_frame.paragraphs[0].text = "publish JSON"
    lbl1.text_frame.paragraphs[0].font.size = Pt(9)
    lbl1.text_frame.paragraphs[0].font.color.rgb = MUTED
    lbl1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.95), Inches(9), Inches(0.8))
    note.text_frame.word_wrap = True
    np = note.text_frame.paragraphs[0]
    np.text = (
        "Production Vercel never mounts \\\\192.168.190.10\\... — only the agent PC on VPN reads the share."
    )
    np.font.size = Pt(12)
    np.font.italic = True
    np.font.color.rgb = MUTED


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Everde AI Operations Portal",
        "Overview & security  •  Jason & Aaron  •  May 2026",
    )
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(8.5), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "https://everde-operations.vercel.app"
    sp.font.size = Pt(14)
    sp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    add_content_slide(
        prs,
        "Agenda",
        [
            "What the portal is and who it is for",
            "How data gets from DataDrops to the browser",
            "Weekly automation (agent PC + schedules)",
            "Security controls in place today",
            "AI assistant — what is sent, who can use it",
            "Gaps, roadmap, and Q&A",
        ],
    )

    add_content_slide(
        prs,
        "What is the portal?",
        [
            "Internal web app for executive and operations dashboards",
            "Freight, Sales Plan, Production & Demand, Retail, Weather, and more",
            "Works on phone, tablet, and desktop — responsive layout",
            "Not Excel on the web: shows pre-calculated metrics from weekly files",
            "Sign-in required with @everde.com Microsoft work accounts (production)",
        ],
    )

    add_architecture_slide(prs)

    add_table_slide(
        prs,
        "Typical weekly / daily automation",
        ["When", "Task", "Output for portal"],
        [
            ["Daily ~8:00 AM", "Sales Plan check", "NOR CAL JSON → Azure Blob"],
            ["Daily ~9:30 AM", "Weather pipeline", "Weather JSON → Azure Blob"],
            ["Monday ~9:00 AM", "Freight rebuild", "Dashboard JSON → Azure Blob"],
            ["Monday ~10:00 AM", "Retail build + extract", "Retail JSON → Azure Blob"],
            ["Monday ~1:30 PM", "Inventory Metrics", "Nursery HTML → git → Vercel"],
        ],
    )

    add_content_slide(
        prs,
        "What users experience",
        [
            "1. Open https://everde-operations.vercel.app",
            "2. Sign in with Microsoft (Everde tenant)",
            "3. Pick a section in the sidebar (Freight, Retail, Sales Plan, …)",
            "4. Dashboard loads from JSON APIs — charts and tables update from Blob",
            "5. Optional: AI analyst in the header (OpenAI or Claude) for questions",
            "Raw Excel files stay on the LAN share; browsers never open UNC paths directly",
        ],
    )

    add_content_slide(
        prs,
        "Security — design principle",
        [
            "Weekly Excel stays on our network.",
            "A trusted PC builds summaries and uploads JSON to Azure.",
            "The portal on Vercel only serves that data to signed-in @everde.com users.",
            "The public internet never sees the file share directly.",
            "",
            "Share access (VPN + permissions) is a separate control from portal access.",
        ],
    )

    add_table_slide(
        prs,
        "Who can access the portal",
        ["Control", "What it does"],
        [
            ["Microsoft Entra ID", "Company sign-in; tokens from Everde tenant"],
            ["@everde.com domain", "Only Everde email identities accepted"],
            ["PORTAL_REQUIRE_AUTH", "All pages and data APIs require sign-in"],
            ["HttpOnly session cookie", "Server-signed session after login — no stored passwords"],
            ["PORTAL_SESSION_SECRET", "Random 32+ char secret on Vercel only (not in git)"],
        ],
    )

    add_table_slide(
        prs,
        "How API and dashboard data are protected",
        ["Control", "What it does"],
        [
            ["Middleware + API guards", "Block unauthenticated access to routes and /api/*"],
            ["No UNC on Vercel", "Production does not read \\\\192.168.190.10\\... from the cloud"],
            ["Azure Blob boundary", "Internet users get published JSON, not raw .xlsb files"],
            ["Connection string server-only", "Azure keys in Vercel env — never committed to git"],
            [".env.local gitignored", "Local secrets stay on the agent / dev machine"],
        ],
    )

    add_content_slide(
        prs,
        "What stays on the LAN",
        [
            "Raw and rebuilt workbooks (freight .xlsb, inventory metrics, sales plan Excel)",
            "Python scripts and pipeline logs on DataDrops / JS Files",
            "Only the agent PC needs VPN/LAN access to those paths",
            "",
            "Agent PC = trust boundary: scheduled tasks run as a dedicated user with share + secrets",
            "Handoff doc for Aaron: scripts/windows/VM_AGENT_HANDOFF.md",
        ],
    )

    add_content_slide(
        prs,
        "AI assistant (in-portal analyst)",
        [
            "Same portal authentication as dashboards",
            "API keys (OpenAI / Anthropic) live server-side only — never in the browser",
            "Sends compacted JSON summaries + route context — not full multi-GB Excel files",
            "Useful for exec questions across freight, sales plan, retail, weather",
            "Backlog: rate limits, audit logging, optional page-only context mode",
        ],
    )

    add_content_slide(
        prs,
        "Honest gaps & talking points",
        [
            "Blob storage: protect connection string like a production password (Azure RBAC)",
            "Share security is separate — VPN users with share access can still open DataDrops",
            "One agent machine should run scheduled tasks (avoid duplicate publishes)",
            "Some nav sections still placeholders or pending data (e.g. Oregon sales plan)",
            "Nursery still deploys via git push; long-term: Blob like freight",
            "Roadmap: full ETL to database, stricter per-section RBAC",
        ],
    )

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, FOREST_DEEP)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Docs: scripts/windows/WEEKLY_DROP_AGENT.md  •  VM_AGENT_HANDOFF.md"
    p2.font.size = Pt(14)
    p2.font.color.rgb = GOLD
    p2.space_before = Pt(20)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
